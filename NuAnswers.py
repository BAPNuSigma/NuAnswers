import streamlit as st
from openai import OpenAI
import pandas as pd
import datetime
import os
import tempfile
from pathlib import Path
import PyPDF2
import docx
import pptx
import csv
import xlrd
import openpyxl
import io
import base64

# Set page config
st.set_page_config(
    page_title="NuAnswers",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide all default Streamlit elements we don't want to show
st.markdown("""
    <style>
        /* Hide page names in sidebar */
        span.css-10trblm.e16nr0p30 {
            display: none;
        }
        /* Hide the default Streamlit menu button */
        button.css-1rs6os.edgvbvh3 {
            display: none;
        }
        /* Hide "streamlit app" text */
        .css-17ziqus {
            display: none;
        }
        /* Hide development mode indicator */
        .stDeployButton {
            display: none;
        }
        /* Sidebar styling */
        .css-1d391kg {
            padding-top: 1rem;
        }
    </style>
""", unsafe_allow_html=True)

# Tutoring hours configuration
TUTORING_HOURS = {
    "Monday": [("10:30", "12:30")],    # 10:30 AM - 12:30 PM
    "Tuesday": [("17:00", "19:00")],   # 5:00 PM - 7:00 PM
    "Wednesday": [("12:00", "14:00")], # 12:00 PM - 2:00 PM
    "Thursday": [("10:30", "12:30")],  # 10:30 AM - 12:30 PM
    "Friday": [("13:00", "15:00")],    # 1:00 PM - 3:00 PM
}

def is_within_tutoring_hours():
    """Check if current time is within tutoring hours."""
    current_time = datetime.datetime.now()
    current_day = current_time.strftime("%A")
    current_time_str = current_time.strftime("%H:%M")
    
    if current_day not in TUTORING_HOURS:
        return False
        
    for start_time, end_time in TUTORING_HOURS[current_day]:
        if start_time <= current_time_str <= end_time:
            return True
    return False

# Initialize session state for registration and tracking
if "registered" not in st.session_state:
    st.session_state.registered = False
if "start_time" not in st.session_state:
    st.session_state.start_time = None
if "user_data" not in st.session_state:
    st.session_state.user_data = {}
if "registration_data" not in st.session_state:
    st.session_state.registration_data = pd.DataFrame(columns=[
        "timestamp", "full_name", "student_id", "email", "grade", "campus",
        "major", "course_name", "course_id", "professor", "usage_time_minutes"
    ])

# Configure data directory
DATA_DIR = Path("/data" if os.path.exists("/data") else ".")
REGISTRATION_DATA_PATH = DATA_DIR / "registration_data.csv"
FEEDBACK_DATA_PATH = DATA_DIR / "feedback_data.csv"
TOPIC_DATA_PATH = DATA_DIR / "topic_data.csv"
COMPLETION_DATA_PATH = DATA_DIR / "completion_data.csv"

# Create data directory if it doesn't exist
DATA_DIR.mkdir(exist_ok=True)

def save_to_csv(data, filepath):
    """Save data to CSV file with error handling"""
    try:
        # Convert to DataFrame if it's a dictionary
        if isinstance(data, dict):
            df = pd.DataFrame([data])
        elif isinstance(data, list):
            df = pd.DataFrame(data)
        else:
            df = data
            
        # Load existing data if file exists
        if filepath.exists():
            existing_df = pd.read_csv(filepath)
            df = pd.concat([existing_df, df], ignore_index=True)
        
        # Save to CSV
        df.to_csv(filepath, index=False)
    except Exception as e:
        st.error(f"Failed to save data to {filepath}: {str(e)}")

def save_registration(user_data, start_time):
    """Save registration data to CSV"""
    end_time = datetime.datetime.now()
    usage_time = (end_time - start_time).total_seconds() / 60
    
    # Create new registration entry
    new_registration = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "full_name": user_data["full_name"],
        "student_id": user_data["student_id"],
        "email": user_data["email"],
        "grade": user_data["grade"],
        "campus": user_data["campus"],
        "major": user_data["major"],
        "course_name": user_data["course_name"],
        "course_id": user_data["course_id"],
        "professor": user_data["professor"],
        "usage_time_minutes": usage_time
    }
    
    # Save to CSV
    save_to_csv(new_registration, REGISTRATION_DATA_PATH)

# Create a sidebar
with st.sidebar:
    st.title("📚 NuAnswers")

# Initialize session state for form
if "form_major" not in st.session_state:
    st.session_state.form_major = "Accounting"

def update_major():
    """Update the major in session state"""
    st.session_state.form_major = st.session_state.temp_major
    st.rerun()

# Registration form
if not st.session_state.registered:
    st.title("📝 Registration Form")
    st.write("Please complete the registration form to use NuAnswers.")
    
    with st.form("registration_form", clear_on_submit=False):
        full_name = st.text_input("Full Name")
        student_id = st.text_input("FDU Student ID (7 digits)")
        
        # Validate student ID format
        is_valid_student_id = student_id.isdigit() and len(student_id) == 7
        if student_id and not is_valid_student_id:
            st.error("FDU Student ID must be exactly 7 digits.")
            
        email = st.text_input("FDU Student Email (@student.fdu.edu or @fdu.edu)")
        
        # Validate email domain
        is_valid_email = False
        if email:
            email = email.lower()
            is_valid_email = email.endswith('@student.fdu.edu') or email.endswith('@fdu.edu')
            if not is_valid_email:
                st.error("Please use your FDU email address (@student.fdu.edu or @fdu.edu)")
        
        grade = st.selectbox("Grade", ["Freshman", "Sophomore", "Junior", "Senior", "Graduate"])
        campus = st.selectbox("Campus", ["Florham", "Metro", "Vancouver"])
        
        # Major selection
        major = st.selectbox(
            "Major",
            ["Accounting", "Finance", "MIS [Management Information Systems]"]
        )
        
        # General course question
        course_name = st.text_input("Which class are you taking that relates to what you need help in?")
        course_id = st.text_input("Course ID (EX: ACCT_####_##)")
        professor = st.text_input("Professor's Name")
        
        submitted = st.form_submit_button("Submit")
        
        if submitted:
            if not all([full_name, student_id, email, course_id, professor]):
                st.error("Please fill in all required fields.")
            elif not is_valid_student_id:
                st.error("Please enter a valid 7-digit FDU Student ID.")
            elif not is_valid_email:
                st.error("Please enter a valid FDU email address.")
            else:
                # Save user data
                st.session_state.user_data = {
                    "full_name": full_name,
                    "student_id": student_id,
                    "email": email,
                    "grade": grade,
                    "campus": campus,
                    "major": major,
                    "course_name": course_name,
                    "course_id": course_id,
                    "professor": professor
                }
                st.session_state.start_time = datetime.datetime.now()
                
                # Save registration data
                save_registration(st.session_state.user_data, st.session_state.start_time)
                
                # Set registered state
                st.session_state.registered = True
                st.rerun()

# Function to extract text from different file types
def extract_text_from_file(file):
    file_extension = Path(file.name).suffix.lower()
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
        tmp_file.write(file.getvalue())
        tmp_file_path = tmp_file.name
    
    try:
        if file_extension == '.pdf':
            text = extract_text_from_pdf(tmp_file_path)
        elif file_extension == '.docx':
            text = extract_text_from_docx(tmp_file_path)
        elif file_extension == '.txt':
            with open(tmp_file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(tmp_file_path)
        elif file_extension == '.csv':
            text = extract_text_from_csv(tmp_file_path)
        elif file_extension in ['.xls', '.xlsx']:
            text = extract_text_from_excel(tmp_file_path)
        else:
            st.error(f"Unsupported file type: {file_extension}")
            return None
    except Exception as e:
        st.error(f"Error processing file: {str(e)}")
        return None
    finally:
        os.unlink(tmp_file_path)
    
    return text

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_csv(file_path):
    text = ""
    with open(file_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            text += ", ".join(row) + "\n"
    return text

def extract_text_from_excel(file_path):
    text = ""
    try:
        # Read all sheets
        excel_file = pd.ExcelFile(file_path)
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            text += f"\nSheet: {sheet_name}\n"
            text += df.to_string(index=False) + "\n"
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return None
    return text

# Initialize session state for uploaded documents and search
if "uploaded_documents" not in st.session_state:
    st.session_state.uploaded_documents = []
if "search_query" not in st.session_state:
    st.session_state.search_query = ""
if "doc_to_delete" not in st.session_state:
    st.session_state.doc_to_delete = None

# Function to search within documents
def search_in_documents(query, documents):
    if not query:
        return documents
    query = query.lower()
    results = []
    for doc in documents:
        if query in doc['name'].lower() or query in doc['content'].lower():
            results.append(doc)
    return results

# Show title and description only after registration
if st.session_state.registered:
    st.title("💬 NuAnswers")
    
    # Check if current time is within tutoring hours
    if is_within_tutoring_hours():
        st.warning("""
        ⚠️ In-person tutoring is currently available! 
        
        Please visit the in-person tutoring session instead of using the bot. 
        The bot will be available after the tutoring session ends.
        """)
        st.stop()

    # Get the API key from environment variable or secrets
    openai_api_key = os.environ.get("OPENAI_API_KEY") or st.secrets.get("OPENAI_API_KEY")

    if not openai_api_key:
        st.error("""
        ⚠️ OpenAI API Key not configured!
        
        To use NuAnswers, you need to:
        1. Get an API key from OpenAI (https://platform.openai.com/api-keys)
        2. Add it as an environment variable:
           - Key: OPENAI_API_KEY
           - Value: your-api-key-here
        3. Restart the application
        
        If you're deploying on Render:
        1. Go to your Render dashboard
        2. Select your service
        3. Go to the "Environment" tab
        4. Add a new environment variable:
           - Key: OPENAI_API_KEY
           - Value: your-api-key-here
        5. Redeploy your service
        """)
        st.stop()

    # Create an OpenAI client
    client = OpenAI(api_key=openai_api_key)

    # File upload section
    st.subheader("📄 Upload Course Materials")
    uploaded_files = st.file_uploader(
        "Upload your course materials (PDF, DOCX, TXT, PPTX, CSV, XLS, XLSX, PNG, JPG, JPEG)",
        type=['pdf', 'docx', 'txt', 'pptx', 'csv', 'xls', 'xlsx', 'png', 'jpg', 'jpeg'],
        accept_multiple_files=True
    )
    
    if uploaded_files:
        for file in uploaded_files:
            if file not in [doc['file'] for doc in st.session_state.uploaded_documents]:
                file_extension = Path(file.name).suffix.lower()
                
                # Handle image files differently
                if file_extension in ['.png', '.jpg', '.jpeg']:
                    # Analyze image content
                    image_analysis = analyze_image(file)
                    
                    st.session_state.uploaded_documents.append({
                        'file': file,
                        'name': file.name,
                        'content': f"[Image Analysis: {image_analysis}]" if image_analysis else f"[Image File: {file.name}]",
                        'is_image': True,
                        'image_analysis': image_analysis
                    })
                    st.success(f"Successfully uploaded and analyzed image {file.name}")
                else:
                    text = extract_text_from_file(file)
                    if text:
                        st.session_state.uploaded_documents.append({
                            'file': file,
                            'name': file.name,
                            'content': text,
                            'is_image': False
                        })
                        st.success(f"Successfully processed {file.name}")
    
    # Search and document management section
    if st.session_state.uploaded_documents:
        st.subheader("📚 Your Uploaded Materials")
        
        # Search bar
        search_col, reorder_col = st.columns([3, 1])
        with search_col:
            st.session_state.search_query = st.text_input("🔍 Search in documents", 
                                                        value=st.session_state.search_query,
                                                        placeholder="Search by filename or content")
        
        # Reorder button
        with reorder_col:
            if st.button("🔄 Reorder Documents"):
                st.session_state.show_reorder = not getattr(st.session_state, 'show_reorder', False)
        
        # Reorder interface
        if getattr(st.session_state, 'show_reorder', False):
            st.info("Drag and drop documents to reorder them")
            for i, doc in enumerate(st.session_state.uploaded_documents):
                cols = st.columns([1, 4, 1])
                with cols[0]:
                    st.write(f"{i+1}.")
                with cols[1]:
                    st.write(doc['name'])
                with cols[2]:
                    if st.button("↑", key=f"up_{i}") and i > 0:
                        st.session_state.uploaded_documents[i], st.session_state.uploaded_documents[i-1] = \
                            st.session_state.uploaded_documents[i-1], st.session_state.uploaded_documents[i]
                        st.rerun()
                    if st.button("↓", key=f"down_{i}") and i < len(st.session_state.uploaded_documents) - 1:
                        st.session_state.uploaded_documents[i], st.session_state.uploaded_documents[i+1] = \
                            st.session_state.uploaded_documents[i+1], st.session_state.uploaded_documents[i]
                        st.rerun()
        
        # Display filtered documents
        filtered_docs = search_in_documents(st.session_state.search_query, st.session_state.uploaded_documents)
        
        if not filtered_docs:
            st.info("No documents match your search query.")
        else:
            for i, doc in enumerate(filtered_docs):
                cols = st.columns([4, 1])
                with cols[0].expander(doc['name']):
                    if doc.get('is_image', False):
                        st.image(doc['file'], caption=doc['name'])
                        if doc.get('image_analysis'):
                            st.markdown("**Image Analysis:**")
                            st.markdown(doc['image_analysis'])
                    else:
                        # Highlight search terms in content
                        content = doc['content']
                        if st.session_state.search_query:
                            query = st.session_state.search_query.lower()
                            start = content.lower().find(query)
                            if start != -1:
                                end = start + len(query)
                                highlighted = (
                                    content[:start] +
                                    f"**{content[start:end]}**" +
                                    content[end:]
                                )
                                st.markdown(highlighted)
                            else:
                                st.text(content[:500] + "..." if len(content) > 500 else content)
                        else:
                            st.text(content[:500] + "..." if len(content) > 500 else content)
                
                # Delete button with confirmation
                if cols[1].button("🗑️", key=f"delete_{i}"):
                    st.session_state.doc_to_delete = doc
                
                # Confirmation dialog
                if st.session_state.doc_to_delete == doc:
                    st.warning(f"Are you sure you want to delete {doc['name']}?")
                    confirm_cols = st.columns(2)
                    if confirm_cols[0].button("Yes, delete it", key=f"confirm_delete_{i}"):
                        st.session_state.uploaded_documents.remove(doc)
                        st.session_state.doc_to_delete = None
                        st.rerun()
                    if confirm_cols[1].button("Cancel", key=f"cancel_delete_{i}"):
                        st.session_state.doc_to_delete = None
                        st.rerun()

    # Create a session state variable to store the chat messages
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "assistant", "content": "Hello! I'm NuAnswers. I'm here to help you understand concepts and work through problems. What would you like to work on today?"}
        ]

    st.write(
        "Hello! I am NuAnswers, Beta Alpha Psi: Nu Sigma Chapter's AI Tutor Bot. I'm here to help you understand concepts and work through problems. "
        "Remember, I won't give you direct answers, but I'll guide you to find them yourself. "
        "I can help you with accounting equations, financial ratios, financial statements, and time value of money concepts."
    )

    # Display the existing chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    # Create a chat input field
    if prompt := st.chat_input("What would you like to work on today?"):
        # Store and display the current prompt
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        # Prepare context from uploaded documents
        context = ""
        if st.session_state.uploaded_documents:
            context = "\n\n".join([f"Document: {doc['name']}\nContent: {doc['content']}" 
                                 for doc in st.session_state.uploaded_documents])
            context = f"Here is the context from uploaded documents:\n\n{context}\n\n"

        # Generate a response using the OpenAI API
        stream = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": f"""You are an Accounting & Finance Tutor. Your role is to guide students through their homework and exam preparation through a conversational, step-by-step approach.

IMPORTANT RULES:
1. NEVER give direct answers or solutions
2. Ask ONE question at a time and wait for the student's response
3. After each student response, ask a follow-up question to guide their thinking
4. If the student's answer is incorrect, ask a guiding question to help them think differently
5. If the student asks for the answer, respond with a question that helps them think about the problem differently
6. Use simple, clear questions that build on each other
7. Focus on one concept or step at a time
8. Validate their understanding before moving to the next step
9. Use encouraging phrases like "Good thinking!" or "You're on the right track!"
10. If the student seems stuck, ask a simpler question that breaks down the problem
11. Use the context from uploaded documents to provide more relevant guidance

Example of good tutoring:
Student: "How do I solve this problem?"
Tutor: "Let's start with the first step. What information do we have in the problem?"
Student: [responds]
Tutor: "Good! Now, what do you think we should do with this information?"
[continue with one question at a time]

Example of bad tutoring:
"Here's how to solve it: First, do this, then do that, then calculate this..."
[giving multiple steps at once]"""},
                {"role": "system", "content": context},
                *[{"role": m["role"], "content": m["content"]} for m in st.session_state.messages]
            ],
            stream=True,
        )

        # Stream the response
        with st.chat_message("assistant"):
            response = st.write_stream(stream)
        st.session_state.messages.append({"role": "assistant", "content": response})

    # Add a logout button
    if st.button("Logout"):
        # Save final usage data before logout
        save_registration(st.session_state.user_data, st.session_state.start_time)
        
        # Reset session state
        st.session_state.registered = False
        st.session_state.start_time = None
        st.session_state.user_data = {}
        st.session_state.messages = []
        st.rerun()

    # Add feedback section after chat
    if st.session_state.messages and len(st.session_state.messages) > 0:
        st.divider()
        st.subheader("📝 Session Feedback")
        
        feedback_col1, feedback_col2, feedback_col3 = st.columns(3)
        
        with feedback_col1:
            topic = st.text_input("What topic did you discuss?", key="feedback_topic")
        
        with feedback_col2:
            rating = st.slider("How helpful was this session?", 1, 5, 3, key="feedback_rating")
        
        with feedback_col3:
            difficulty = st.slider("How difficult was the topic?", 1, 5, 3, key="feedback_difficulty")
        
        if st.button("Submit Feedback"):
            if topic:
                save_feedback(rating, topic, difficulty)
                track_topic(topic, difficulty)
                track_completion(True)
                st.success("Thank you for your feedback!")
            else:
                st.warning("Please enter the topic discussed.")

def show_admin_panel():
    """Display admin panel with registration statistics"""
    st.header("Admin Panel")
    
    try:
        # Load data from CSV
        csv_path = "registration_data.csv"
        if os.path.exists(csv_path):
            df = pd.read_csv(csv_path)
            
            # Convert timestamp to datetime
            df['timestamp'] = pd.to_datetime(df['timestamp'])
            
            # Calculate statistics
            total_registrations = len(df)
            total_usage_minutes = df['usage_time_minutes'].sum()
            avg_usage_minutes = df['usage_time_minutes'].mean()
            
            # Display statistics
            st.metric("Total Registrations", total_registrations)
            st.metric("Total Usage Time (minutes)", f"{total_usage_minutes:.2f}")
            st.metric("Average Usage Time (minutes)", f"{avg_usage_minutes:.2f}")
            
            # Show daily statistics
            st.subheader("Daily Statistics")
            daily_stats = df.groupby(df['timestamp'].dt.date).agg({
                'student_id': 'count',
                'usage_time_minutes': ['sum', 'mean']
            }).reset_index()
            daily_stats.columns = ['Date', 'Registrations', 'Total Minutes', 'Avg Minutes']
            st.dataframe(daily_stats.sort_values('Date', ascending=False))
            
            # Show raw data
            st.subheader("Raw Registration Data")
            st.dataframe(df.sort_values('timestamp', ascending=False))
        else:
            st.info("No registration data available yet.")
    except Exception as e:
        st.error(f"Error loading registration data: {str(e)}")

# Initialize session state for feedback and topic tracking
if "feedback_data" not in st.session_state:
    st.session_state.feedback_data = []
if "topic_data" not in st.session_state:
    st.session_state.topic_data = []
if "completion_data" not in st.session_state:
    st.session_state.completion_data = []

def save_feedback(rating, topic, difficulty):
    """Save feedback data"""
    feedback_entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "student_id": st.session_state.user_data.get("student_id"),
        "course_id": st.session_state.user_data.get("course_id"),
        "rating": rating,
        "topic": topic,
        "difficulty": difficulty
    }
    st.session_state.feedback_data.append(feedback_entry)
    save_to_csv(feedback_entry, FEEDBACK_DATA_PATH)

def track_topic(topic, difficulty=None):
    """Track topic data"""
    topic_entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "student_id": st.session_state.user_data.get("student_id"),
        "course_id": st.session_state.user_data.get("course_id"),
        "topic": topic,
        "difficulty": difficulty
    }
    st.session_state.topic_data.append(topic_entry)
    save_to_csv(topic_entry, TOPIC_DATA_PATH)

def track_completion(completed):
    """Track course completion"""
    completion_entry = {
        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "student_id": st.session_state.user_data.get("student_id"),
        "course_id": st.session_state.user_data.get("course_id"),
        "completed": completed
    }
    st.session_state.completion_data.append(completion_entry)
    save_to_csv(completion_entry, COMPLETION_DATA_PATH)

def encode_image_to_base64(file):
    """Convert uploaded image file to base64 string"""
    return base64.b64encode(file.getvalue()).decode('utf-8')

def analyze_image(file):
    """Analyze image content using OpenAI's GPT-4 Vision model"""
    try:
        base64_image = encode_image_to_base64(file)
        
        response = client.chat.completions.create(
            model="gpt-4-vision-preview",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Please analyze this image in the context of accounting, finance, or business studies. Describe any relevant equations, problems, charts, or concepts shown."},
                        {
                            "type": "image_url",
                            "image_url": f"data:image/jpeg;base64,{base64_image}"
                        }
                    ]
                }
            ],
            max_tokens=300
        )
        
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Error analyzing image: {str(e)}")
        return None
